#!/usr/bin/env python3
from pathlib import Path

import matplotlib.pyplot as plt
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
CHUNKS = [2000, 8000, 32000, 128000, 512000]
TOTAL_US = [337844, 91303, 45748, 47482, 44046]
MEDIAN_US = [345594, 85771, 39882, 28276, 25782]
P90_US = [391651, 105447, 61045, 115303, 110104]
CV = [0.14, 0.16, 0.26, 0.74, 0.80]


def rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        run.font.name = "Aptos Display"
    return p


def add_para(doc, text, bold_prefix=None):
    p = doc.add_paragraph()
    if bold_prefix and text.startswith(bold_prefix):
        r = p.add_run(bold_prefix)
        r.bold = True
        p.add_run(text[len(bold_prefix):])
    else:
        p.add_run(text)
    return p


def make_docx():
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.7)
    sec.bottom_margin = Inches(0.7)
    sec.left_margin = Inches(0.8)
    sec.right_margin = Inches(0.8)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("Mini 2: Distributed 311 Query Chunking")
    r.bold = True
    r.font.size = Pt(20)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("Anukrithi Myadala, Asim Mohammed").italic = True

    add_heading(doc, "Research Question", 1)
    add_para(doc, "How does chunk size affect a distributed, sharded 311 result set when A is the only client-facing process?")

    add_heading(doc, "Mini 1 Takeaways Applied", 1)
    for item in [
        "Use typed fields instead of strings for everything.",
        "Avoid shared merge contention by gathering child replies into separate buffers.",
        "Reserve or reuse storage where possible.",
        "Report failures and limitations, not only successful timings.",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    add_heading(doc, "Design", 1)
    add_para(doc, "The cluster uses a configured scatter-gather tree: A -> B,H,G,I; B -> C,D,E; E -> F. A is the only public entry point. B-I own typed binary shards, and I is implemented in Python.")
    add_para(doc, "The data source is NYC 311 Service Requests from 2020 to Present. We used a 90,000-row subset from the public CSV so the test was repeatable and the Canvas submission did not include the large dataset.")
    add_para(doc, "Each 311 row is stored as a compact 20-byte record: unique key, latitude, longitude, incident zip, created year, status code, and borough code.")
    add_para(doc, "The design uses unary gRPC calls with explicit offsets. A streaming prototype was dropped because the assignment required non-streaming gRPC and the explicit offsets made failure behavior easier to test.")

    add_heading(doc, "Course and Lab Ideas Used", 1)
    for item in [
        "The basic-grpc lab shaped the protobuf/service structure and unary call pattern.",
        "The leader-adv lab shaped A as the coordinator while B-I do shard work.",
        "The socket interoperability lab made us treat the C++/Python 20-byte record layout as a contract.",
        "Socket and messaging lectures motivated the chunk-size experiment.",
        "The sharding lecture motivated splitting the binary data across B-I.",
        "MPI round/baton-style labs influenced the fairness test.",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    add_heading(doc, "Mini 2 Questions Answered", 1)
    questions = [
        ("Performance/resources", "30-run chunk sweep measured total time, chunk count, avg RPC time, min RPC time, and max RPC time."),
        ("Conserve memory", "20-byte typed records, explicit chunk offsets, and a 1 MB max chunk instead of unbounded responses."),
        ("Fairness/balance", "Four-client run showed equal chunk turns, but not identical finish times."),
        ("Flexible overlay", "Host/process/tree settings live in YAML; identity and config path are runtime arguments."),
        ("No flat shortcut", "Final tree uses A -> B,H,G,I; B -> C,D,E; E -> F."),
        ("Request failure/abandonment", "Fail-fast before gather; cache-complete behavior after gather; no speculative prefetching because it would increase A memory pressure."),
    ]
    table_q = doc.add_table(rows=1, cols=2)
    table_q.style = "Table Grid"
    table_q.rows[0].cells[0].text = "Prompt challenge"
    table_q.rows[0].cells[1].text = "Our answer"
    for left, right in questions:
        cells = table_q.add_row().cells
        cells[0].text = left
        cells[1].text = right

    add_heading(doc, "Measurement Plan", 1)
    add_para(doc, "The course notes recommend 15-30 runs to form an average and discard clear outliers. The final table below uses 30 runs per chunk size on two laptops.")

    add_heading(doc, "Two-Host Results", 1)
    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Chunk bytes"
    hdr[1].text = "Avg total us"
    hdr[2].text = "Avg chunks"
    hdr[3].text = "Avg RPC us"
    rows = [
        (2000, 337844, 800, 422),
        (8000, 91303, 200, 456),
        (32000, 45748, 50, 914),
        (128000, 47482, 13, 3652),
        (512000, 44046, 4, 11011),
    ]
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = str(value)

    add_para(doc, "Result: 512 KB chunks completed the same 80,000-record response about 7.7x faster than 2 KB chunks across two laptops because the client needed 4 chunks instead of 800. The more interesting result is the knee in the curve: 32 KB used 50 chunks and averaged 45.7 ms, while 512 KB used only 4 chunks and averaged 44.0 ms. That 1.7 ms difference means page count stopped being the only important cost after about 32 KB.")
    add_para(doc, "The accurate interpretation is that chunk size shifts cost between repeated paging overhead and a fixed gather/transfer floor. The first page triggers A to gather and cache the full 1.6 MB response from B-I. Later pages come from A's cache, so small chunks pay many local unary RPC/cache-copy steps. Once chunks are large enough, the remaining time is dominated by the one-time gather, payload movement, serialization/copying, and Wi-Fi tail spikes.")

    add_heading(doc, "Mean, Median, and Tail Behavior", 1)
    tail_table = doc.add_table(rows=1, cols=5)
    tail_table.style = "Table Grid"
    hdr = tail_table.rows[0].cells
    hdr[0].text = "Chunk bytes"
    hdr[1].text = "Mean ms"
    hdr[2].text = "Median ms"
    hdr[3].text = "P90 ms"
    hdr[4].text = "CV"
    rows_tail = [
        (2000, 337.8, 345.6, 391.7, 0.14),
        (8000, 91.3, 85.8, 105.4, 0.16),
        (32000, 45.7, 39.9, 61.0, 0.26),
        (128000, 47.5, 28.3, 115.3, 0.74),
        (512000, 44.0, 25.8, 110.1, 0.80),
    ]
    for row in rows_tail:
        cells = tail_table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = str(value)
    add_para(doc, "CV is the sample standard deviation divided by the mean across the 30 runs.")
    add_para(doc, "The medians show that large chunks were usually faster, but the high P90/CV values show why their means were not much better than 32 KB. With only 4-13 chunks, one slow remote gather or large response spike can dominate an entire run. A 32 KB chunk is therefore a good robust operating point for this setup: near-fastest mean, lower tail risk, and the fairness test was run at this size.")

    add_heading(doc, "Fairness Result", 1)
    add_para(doc, "Four clients at 32 KB chunks each completed 50 chunks. The slowest client was about 4.3% slower than the fastest, so the queue balances chunk turns but does not guarantee identical wall-clock finish times.")

    add_heading(doc, "Failures Found and Fixed", 1)
    failures = [
        "Early per-RPC channel/stub setup made the prototype measure setup overhead too much. Child stubs are now created once during startup.",
        "A streaming/async detour was abandoned because the final design needed unary gRPC and clearer failure behavior.",
        "Old processes were already bound to ports 50051, 50052, and 50058, causing the client to hit the wrong server.",
        "The first tree derivation only contacted B's subtree. The directed tree is now explicit in nodes.yaml.",
        "Partial child failures could be cached. Child fetch failures now fail the request.",
        "Client request ids were reused. They now include a per-run timestamp.",
        "Python and C++ initially had to be checked carefully for exact binary record size and field order. The final record is a validated 20-byte layout.",
        "The Python node failed under system Python without yaml. The launcher now uses venv/bin/python when available.",
        "Host2 Homebrew Python loaded macOS's older libexpat, which broke ensurepip. We used Python 3.12 with Homebrew's expat path.",
        "Node I initially used fallback sample data until we copied the real shards to host2 and restarted the node.",
        "The first benchmark included one-record chunks, which was too slow for normal iteration. Tiny chunks are now opt-in.",
        "Large requested chunks were clamped to 64 KB. The cap is now 1 MB.",
        "With H down before a new request, the client fails clearly. If H dies after A has cached the gathered result, that same request can still finish from cache.",
    ]
    for item in failures:
        doc.add_paragraph(item, style="List Bullet")

    add_heading(doc, "Conclusion", 1)
    add_para(doc, "The final result supports one focused conclusion: chunk size has a knee around 32 KB. Below 32 KB, repeated pages dominate. Above 32 KB, the fixed gather/transfer floor and tail variance dominate. If optimizing only for mean or median, 512 KB wins by a small amount. If optimizing for robustness, fairness, and tail behavior, 32 KB is the better operating point for this two-laptop setup. This conclusion only became trustworthy after we fixed port collisions, topology mistakes, cache behavior, the hidden chunk cap, Python setup problems, missing shard deployment, and binary-layout validation.")

    add_heading(doc, "Individual Contributions", 1)
    add_para(doc, "Anukrithi Myadala focused on Mini 1 feedback analysis, runbooks, two-computer setup, result collection, report, and presentation framing. Asim Mohammed contributed to the cluster/protobuf implementation and helped with final code cleanup and validation.")

    add_heading(doc, "References", 1)
    for item in [
        "NYC Open Data / Data.gov, 311 Service Requests from 2020 to Present: https://catalog.data.gov/dataset/311-service-requests-from-2010-to-present",
        "NYC Open Data, 311 Service Requests Updates: https://opendata.cityofnewyork.us/311-service-requests-from-2010-to-present-updates/",
        "gRPC, Performance Best Practices: https://grpc.io/docs/guides/performance/",
        "Protocol Buffers, Encoding: https://protobuf.dev/programming-guides/encoding/",
        "Protocol Buffers, Language Guide (proto3): https://protobuf.dev/programming-guides/proto3/",
        "Python struct documentation for the exact C++/Python binary record contract: https://docs.python.org/3/library/struct.html",
        "CMake add_custom_command documentation for generating protobuf/gRPC sources during the build: https://cmake.org/cmake/help/latest/command/add_custom_command.html",
        "gRPC C++ basics tutorial for protoc/grpc_cpp_plugin build pattern: https://grpc.io/docs/languages/cpp/basics/",
        "PyYAML documentation for loading YAML configuration on the Python node: https://pyyaml.org/wiki/PyYAMLDocumentation",
        "yaml-cpp project documentation for C++ YAML parsing: https://github.com/jbeder/yaml-cpp",
        "python-pptx documentation for generating the one-slide poster: https://python-pptx.readthedocs.io/",
        "Matplotlib barh documentation for the horizontal poster chart: https://matplotlib.org/stable/api/_as_gen/matplotlib.pyplot.barh.html",
        "AMD Vitis HLS Documentation, Data Structure Padding: https://docs.amd.com/r/2024.1-English/ug1399-vitis-hls/Data-Structure-Padding",
        "Course lectures on messaging/socket costs, sharding, parallelism, failure behavior, and benchmarking.",
        "Course labs covering basic gRPC, leader coordination, MPI round/baton behavior, and sockets.",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    doc.save(ROOT / "mini2-report.docx")


def make_poster_chart(path):
    labels = ["2 KB  | 800 pages", "8 KB  | 200 pages", "32 KB | 50 pages", "128 KB | 13 pages", "512 KB | 4 pages"]
    mean_ms = [v / 1000 for v in TOTAL_US]
    median_ms = [v / 1000 for v in MEDIAN_US]
    p90_ms = [v / 1000 for v in P90_US]
    colors = ["#fb7185", "#fbbf24", "#38bdf8", "#34d399", "#a78bfa"]

    plt.figure(figsize=(8.4, 4.9), facecolor="#111827")
    ax = plt.gca()
    ax.set_facecolor("#111827")
    bars = ax.barh(labels, mean_ms, color=colors, height=0.52, alpha=0.95, label="mean")
    ax.invert_yaxis()
    ax.set_xlim(0, 430)
    ax.set_xlabel("total request time, ms", color="#cbd5e1", labelpad=10)
    ax.tick_params(axis="x", colors="#94a3b8", labelsize=10)
    ax.tick_params(axis="y", colors="#f8fafc", labelsize=11.5)
    ax.grid(axis="x", color="#334155", linewidth=0.7, alpha=0.55)
    ax.set_axisbelow(True)
    for spine in ax.spines.values():
        spine.set_visible(False)
    for bar, value, med, p90 in zip(bars, mean_ms, median_ms, p90_ms):
        y = bar.get_y() + bar.get_height() / 2
        ax.plot(med, y, marker="o", markersize=7, color="#f8fafc", markeredgecolor="#111827", markeredgewidth=1.2)
        ax.plot([med, p90], [y, y], color="#e2e8f0", linewidth=1.2, alpha=0.45)
        ax.plot(p90, y, marker="|", markersize=15, color="#e2e8f0", alpha=0.85)
        ax.text(
            bar.get_width() + 7,
            y,
            f"mean {value:.0f}",
            va="center",
            ha="left",
            color="#f8fafc",
            fontsize=12,
            fontweight="bold",
        )
    ax.text(
        0,
        -0.92,
        "Mean bars, median dots, P90 ticks",
        color="#e2e8f0",
        fontsize=13,
        fontweight="bold",
    )
    ax.text(
        0,
        -0.55,
        "Knee near 32 KB: mean flattens after that, but tail risk grows.",
        color="#94a3b8",
        fontsize=10.5,
    )
    plt.tight_layout(pad=1.8)
    plt.savefig(path, dpi=220, facecolor="#111827")
    plt.close()


def make_poster():
    poster_chart = ROOT / "results" / "poster_chunk_chart.png"
    make_poster_chart(poster_chart)

    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#070b16")

    def add_box(x, y, w, h, fill, line="#1e293b", radius=True):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line)
        return shape

    def add_text(x, y, w, h, text, size, color="#f8fafc", bold=False, align=None):
        shape = slide.shapes.add_textbox(x, y, w, h)
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = PptInches(0.04)
        tf.margin_right = PptInches(0.04)
        tf.margin_top = PptInches(0.02)
        tf.margin_bottom = PptInches(0.02)
        p = tf.paragraphs[0]
        p.text = text
        if align:
            p.alignment = align
        p.font.size = PptPt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        return shape

    def add_line(x1, y1, x2, y2, color="#64748b", width=1.2):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = rgb(color)
        line.line.width = PptPt(width)
        return line

    def add_node(cx, cy, label, fill="#172033", line="#64748b", color="#f8fafc"):
        size = PptInches(0.36)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            cx - size / 2,
            cy - size / 2,
            size,
            size,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = PptPt(10)
        p.font.color.rgb = rgb(color)
        return shape

    add_box(PptInches(0), PptInches(0), PptInches(13.333), PptInches(0.10), "#f59e0b", "#f59e0b", False)
    add_text(PptInches(0.32), PptInches(0.24), PptInches(3.55), PptInches(0.42), "THE 32 KB KNEE", 28, "#facc15", True)
    add_text(
        PptInches(3.98),
        PptInches(0.28),
        PptInches(7.5),
        PptInches(0.34),
        "Why 50 pages at 32 KB almost tied 4 pages at 512 KB",
        14.5,
        "#e2e8f0",
        True,
    )
    add_text(PptInches(11.6), PptInches(0.34), PptInches(1.3), PptInches(0.24), "Anukrithi + Asim", 9, "#94a3b8", False, PP_ALIGN.RIGHT)
    add_text(
        PptInches(4.12),
        PptInches(0.68),
        PptInches(8.2),
        PptInches(0.22),
        "NYC 311 2020-present | 80k returned rows | 20-byte records | A-F host1, G/H/I host2 | 30 runs/chunk",
        8.8,
        "#94a3b8",
    )

    add_box(PptInches(0.3), PptInches(0.98), PptInches(5.22), PptInches(3.5), "#111827", "#334155")
    slide.shapes.add_picture(str(poster_chart), PptInches(0.48), PptInches(1.17), width=PptInches(4.88), height=PptInches(3.08))

    add_box(PptInches(5.72), PptInches(0.98), PptInches(3.25), PptInches(3.5), "#0f172a", "#334155")
    add_text(PptInches(5.92), PptInches(1.12), PptInches(2.85), PptInches(0.24), "Scatter-gather fan-in", 12.5, "#f8fafc", True, PP_ALIGN.CENTER)
    add_text(PptInches(6.02), PptInches(1.38), PptInches(2.65), PptInches(0.2), "Client only calls A", 9.2, "#38bdf8", True, PP_ALIGN.CENTER)
    pts = {
        "Client": (PptInches(6.15), PptInches(1.98)),
        "A": (PptInches(6.95), PptInches(1.98)),
        "B": (PptInches(6.25), PptInches(2.65)),
        "H": (PptInches(6.85), PptInches(2.65)),
        "G": (PptInches(7.45), PptInches(2.65)),
        "I": (PptInches(8.05), PptInches(2.65)),
        "C": (PptInches(5.94), PptInches(3.23)),
        "D": (PptInches(6.32), PptInches(3.23)),
        "E": (PptInches(6.70), PptInches(3.23)),
        "F": (PptInches(6.70), PptInches(3.72)),
    }
    add_line(*pts["Client"], *pts["A"], "#38bdf8", 2.0)
    for child in ["B", "H", "G", "I"]:
        add_line(*pts["A"], *pts[child], "#64748b", 1.1)
    for child in ["C", "D", "E"]:
        add_line(*pts["B"], *pts[child], "#64748b", 1.1)
    add_line(*pts["E"], *pts["F"], "#64748b", 1.1)
    add_box(pts["Client"][0] - PptInches(0.32), pts["Client"][1] - PptInches(0.16), PptInches(0.64), PptInches(0.32), "#1e293b", "#38bdf8")
    add_text(pts["Client"][0] - PptInches(0.29), pts["Client"][1] - PptInches(0.09), PptInches(0.58), PptInches(0.16), "Client", 7.8, "#38bdf8", True, PP_ALIGN.CENTER)
    add_node(*pts["A"], "A", "#3b2f12", "#facc15", "#facc15")
    for n in ["B", "C", "D", "E", "F", "G", "H"]:
        add_node(*pts[n], n)
    add_node(*pts["I"], "I", "#13251f", "#34d399", "#bbf7d0")
    add_text(PptInches(7.45), PptInches(1.78), PptInches(1.2), PptInches(0.34), "first page gathers\nB-I once", 7.8, "#94a3b8", False, PP_ALIGN.CENTER)
    add_text(PptInches(5.98), PptInches(4.05), PptInches(2.7), PptInches(0.22), "A materializes/cache-builds the 1.6 MB result", 7.8, "#cbd5e1", False, PP_ALIGN.CENTER)

    add_box(PptInches(9.15), PptInches(0.98), PptInches(3.9), PptInches(3.5), "#101826", "#334155")
    add_text(PptInches(9.38), PptInches(1.12), PptInches(3.45), PptInches(0.24), "Total time components", 12.5, "#f8fafc", True, PP_ALIGN.CENTER)
    pieces = [
        ("scatter-gather", "B-I fan-in over Wi-Fi", "#38bdf8"),
        ("cache materialization", "A assembles 1.6 MB", "#34d399"),
        ("amortized page cost", "N QueryOnce calls", "#facc15"),
        ("serialization/copy", "protobuf bytes + slices", "#a78bfa"),
    ]
    x = PptInches(9.42)
    y = PptInches(1.58)
    for i, (top, bottom, color) in enumerate(pieces):
        add_box(x, y + PptInches(0.42 * i), PptInches(3.22), PptInches(0.33), "#172033", color)
        add_text(x + PptInches(0.08), y + PptInches(0.42 * i) + PptInches(0.06), PptInches(1.28), PptInches(0.16), top, 7.7, color, True)
        add_text(x + PptInches(1.42), y + PptInches(0.42 * i) + PptInches(0.06), PptInches(1.7), PptInches(0.16), bottom, 7.1, "#cbd5e1")
        if i < len(pieces) - 1:
            add_text(x + PptInches(1.52), y + PptInches(0.42 * i) + PptInches(0.29), PptInches(0.16), PptInches(0.14), "+", 8.5, "#94a3b8", True, PP_ALIGN.CENTER)
    add_box(PptInches(9.48), PptInches(3.48), PptInches(3.12), PptInches(0.55), "#271c13", "#f59e0b")
    add_text(PptInches(9.62), PptInches(3.58), PptInches(2.82), PptInches(0.28), "32 KB has mostly paid off the page-count penalty.", 9.1, "#fde68a", True, PP_ALIGN.CENTER)

    add_box(PptInches(0.3), PptInches(4.72), PptInches(3.35), PptInches(2.16), "#111827", "#334155")
    add_text(PptInches(0.5), PptInches(4.86), PptInches(2.95), PptInches(0.24), "Page overhead multiplies", 12.0, "#f8fafc", True)
    rows = [
        ("2 KB", "800 pages", 18, "#fb7185"),
        ("32 KB", "50 pages", 9, "#38bdf8"),
        ("512 KB", "4 pages", 3, "#a78bfa"),
    ]
    for r, (label, pages, blocks, color) in enumerate(rows):
        yy = PptInches(5.26 + 0.42 * r)
        add_text(PptInches(0.52), yy, PptInches(0.62), PptInches(0.18), label, 8.9, color, True)
        add_text(PptInches(1.12), yy, PptInches(0.7), PptInches(0.18), pages, 7.9, "#cbd5e1")
        for b in range(blocks):
            bx = PptInches(1.9 + 0.06 * b)
            add_box(bx, yy + PptInches(0.01), PptInches(0.045), PptInches(0.16), color, color, False)
    add_text(PptInches(0.52), PptInches(6.55), PptInches(2.9), PptInches(0.18), "338 ms is all 800 pages, not one chunk.", 8.0, "#facc15", True, PP_ALIGN.CENTER)

    add_box(PptInches(3.87), PptInches(4.72), PptInches(4.35), PptInches(2.16), "#0f172a", "#334155")
    add_text(PptInches(4.08), PptInches(4.86), PptInches(3.95), PptInches(0.24), "Fairness: equal turns, not identical latency", 12.0, "#f8fafc", True, PP_ALIGN.CENTER)
    colors_clients = ["#38bdf8", "#34d399", "#facc15", "#a78bfa"]
    for i, (client, color) in enumerate(zip(["cli1", "cli2", "cli3", "cli4"], colors_clients)):
        yy = PptInches(5.28 + 0.28 * i)
        add_text(PptInches(4.12), yy, PptInches(0.45), PptInches(0.16), client, 7.5, color, True)
        add_line(PptInches(4.6), yy + PptInches(0.08), PptInches(5.28), yy + PptInches(0.08), color, 1.2)
    add_box(PptInches(5.38), PptInches(5.33), PptInches(0.78), PptInches(0.72), "#172033", "#facc15")
    add_text(PptInches(5.46), PptInches(5.49), PptInches(0.62), PptInches(0.16), "A", 13, "#facc15", True, PP_ALIGN.CENTER)
    add_text(PptInches(5.44), PptInches(5.73), PptInches(0.66), PptInches(0.14), "queue", 7, "#facc15", True, PP_ALIGN.CENTER)
    for i, color in enumerate(colors_clients):
        add_line(PptInches(6.2), PptInches(5.68), PptInches(6.7 + 0.26 * i), PptInches(5.68), color, 1.0)
        add_box(PptInches(6.78 + 0.26 * i), PptInches(5.55), PptInches(0.16), PptInches(0.24), color, color, False)
    add_text(PptInches(4.18), PptInches(6.3), PptInches(3.72), PptInches(0.32), "Each client received 50 chunks; finish spread was 4.3% (116.8-121.8 ms).", 8.2, "#cbd5e1", False, PP_ALIGN.CENTER)

    add_box(PptInches(8.45), PptInches(4.72), PptInches(4.6), PptInches(2.16), "#172033", "#475569")
    add_text(PptInches(8.72), PptInches(4.85), PptInches(1.15), PptInches(0.42), "7.7x", 30, "#facc15", True)
    add_text(PptInches(9.85), PptInches(4.96), PptInches(2.65), PptInches(0.28), "338 ms -> 44 ms", 17, "#38bdf8", True)
    add_text(PptInches(8.72), PptInches(5.35), PptInches(3.8), PptInches(0.2), "2 KB to 512 KB total mean", 8.7, "#94a3b8")
    add_box(PptInches(8.76), PptInches(5.72), PptInches(1.6), PptInches(0.62), "#101826", "#38bdf8")
    add_box(PptInches(10.62), PptInches(5.72), PptInches(1.6), PptInches(0.62), "#101826", "#a78bfa")
    add_text(PptInches(8.88), PptInches(5.84), PptInches(1.36), PptInches(0.16), "32 KB", 10.5, "#38bdf8", True, PP_ALIGN.CENTER)
    add_text(PptInches(8.88), PptInches(6.05), PptInches(1.36), PptInches(0.16), "45.7 ms / P90 61", 7.4, "#dbeafe", False, PP_ALIGN.CENTER)
    add_text(PptInches(10.74), PptInches(5.84), PptInches(1.36), PptInches(0.16), "512 KB", 10.5, "#a78bfa", True, PP_ALIGN.CENTER)
    add_text(PptInches(10.74), PptInches(6.05), PptInches(1.36), PptInches(0.16), "44.0 ms / P90 110", 7.4, "#ede9fe", False, PP_ALIGN.CENTER)
    add_text(PptInches(8.82), PptInches(6.5), PptInches(3.75), PptInches(0.18), "At 32 KB, page-count cost is mostly amortized.", 8.4, "#facc15", True, PP_ALIGN.CENTER)

    add_text(PptInches(0.45), PptInches(7.04), PptInches(12.45), PptInches(0.22), "Takeaway: first page gathers/cache-builds once; chunk size controls how many later client-to-A page costs are repeated.", 10.7, "#e2e8f0", True, PP_ALIGN.CENTER)

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = (
        "I want to focus on one finding from our final run, not walk through the "
        "whole project checklist. We used NYC 311 service requests, converted a "
        "90,000 row subset into compact 20 byte records, and returned 80,000 "
        "records through our two laptop, nine node gRPC tree. The only variable "
        "in this graph is chunk size. At first the result looks simple: 2 KB "
        "chunks took about 338 ms, and 512 KB chunks took about 44 ms, so the "
        "extreme speedup is 7.7x. But the part we found more interesting is that "
        "32 KB used 50 pages and 512 KB used only 4 pages, yet their means were "
        "almost tied: 45.7 ms versus 44.0 ms. The slide is explaining why that "
        "happens.\n\n"
        "The tree diagram is the starting point. The client does not directly "
        "call B through I. The client only calls A, because A is the public entry "
        "point and owns the scatter-gather fan-in. On the first page, A has to "
        "query B, H, G, and I. B gathers C, D, and E, and E gathers F. Once those "
        "payloads return, A materializes the full 1.6 MB result and stores it in "
        "the cache under the request id. That first page is special because it "
        "pays the distributed gather cost. Later pages for the same request id do "
        "not gather the whole tree again; A slices the cached payload and returns "
        "the next chunk.\n\n"
        "That is why total time is not only bandwidth. The cost has four pieces: "
        "scatter-gather fan-in over Wi-Fi, cache materialization at A, amortized "
        "page cost for each QueryOnce call, and serialization/copy work for the "
        "protobuf bytes. Small chunks help on one part: each response is tiny. "
        "But they hurt badly on another part: we repeat the page path hundreds of "
        "times. Large chunks reduce the number of page calls, but each response "
        "is larger and the tail latency becomes more visible.\n\n"
        "The bottom-left diagram is there to avoid a common misunderstanding. "
        "The 338 ms at 2 KB is not one chunk response. It is all 800 page "
        "responses combined. Each 2 KB page is cheap, about 0.42 ms on average, "
        "but 800 cheap operations still add up. At 32 KB, the client only needs "
        "50 pages, so most of the repeated paging overhead has already been "
        "amortized. Going from 50 pages down to 4 pages still helps a little, "
        "but now the fixed gather/cache floor and payload movement dominate. "
        "That is why 32 KB and 512 KB have almost the same mean.\n\n"
        "The fairness box shows why we did not just say 512 KB is the answer. "
        "At the 32 KB knee, four clients each got exactly 50 chunks through A's "
        "queue. Their finish times were not identical, but they were close: "
        "116.8 to 121.8 ms, a 4.3 percent spread. So this is opportunity "
        "fairness. A gives clients equal turns, but network timing still creates "
        "slightly different wall-clock completion times.\n\n"
        "The conclusion is that the unique finding is not simply bigger chunks "
        "are faster. Chunk size shifts the bottleneck. Below 32 KB, repeated "
        "client-to-A paging dominates. Above 32 KB, the fixed gather/cache floor, "
        "serialization and copy work, and tail latency dominate. If we optimize "
        "only for mean or median, 512 KB wins by a little. If we care about "
        "robustness and fairness, 32 KB is the better operating point for this "
        "setup because it is nearly as fast but has much lower P90."
    )

    prs.save(ROOT / "mini2-poster.pptx")


def main():
    make_docx()
    make_poster()
    plt.figure(figsize=(8, 4.4))
    plt.bar([str(x) for x in CHUNKS], TOTAL_US, color="#50beaa")
    plt.title("Total Request Time vs. Chunk Size")
    plt.xlabel("Chunk bytes")
    plt.ylabel("Avg total us")
    plt.tight_layout()
    plt.savefig(ROOT / "results" / "chunk_sweep_chart.png", dpi=160)


if __name__ == "__main__":
    main()
